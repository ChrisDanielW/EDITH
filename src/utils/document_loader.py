"""
Document Loader for EDITH
Handles loading and extracting text from various document formats:
- PDF files
- Word documents (.docx, .doc)
- PowerPoint presentations (.pptx)
- Images (with OCR)
- Text files
"""

import os
from typing import List, Dict, Optional
from pathlib import Path
import logging

# Document processing imports
try:
    import pypdf
    from pypdf import PdfReader
except ImportError:
    pypdf = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentLoader:
    """Loads and extracts text from various document formats"""
    
    SUPPORTED_FORMATS = {
        'pdf': ['.pdf'],
        'word': ['.docx', '.doc'],
        'powerpoint': ['.pptx'],
        'image': ['.png', '.jpg', '.jpeg', '.tiff', '.bmp'],
        'text': ['.txt', '.md', '.markdown']
    }
    
    def __init__(self, use_ocr: bool = True):
        """
        Initialize the document loader
        
        Args:
            use_ocr: Whether to use OCR for image-based documents
        """
        self.use_ocr = use_ocr
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check if required dependencies are installed"""
        if pypdf is None:
            logger.warning("pypdf not installed. PDF support limited.")
        if Document is None:
            logger.warning("python-docx not installed. Word document support disabled.")
        if Presentation is None:
            logger.warning("python-pptx not installed. PowerPoint support disabled.")
        if Image is None or pytesseract is None:
            logger.warning("PIL or pytesseract not installed. OCR support disabled.")
            self.use_ocr = False
    
    def load_document(self, file_path: str) -> Dict[str, any]:
        """
        Load a document and extract its text content
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary containing:
                - text: Extracted text content
                - metadata: Document metadata (filename, type, page_count, etc.)
                - success: Boolean indicating if extraction was successful
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return {
                'text': '',
                'metadata': {'filename': file_path.name, 'error': 'File not found'},
                'success': False
            }
        
        file_extension = file_path.suffix.lower()
        
        # Route to appropriate loader
        if file_extension in self.SUPPORTED_FORMATS['pdf']:
            return self._load_pdf(file_path)
        elif file_extension in self.SUPPORTED_FORMATS['word']:
            return self._load_word(file_path)
        elif file_extension in self.SUPPORTED_FORMATS['powerpoint']:
            return self._load_powerpoint(file_path)
        elif file_extension in self.SUPPORTED_FORMATS['image']:
            return self._load_image(file_path)
        elif file_extension in self.SUPPORTED_FORMATS['text']:
            return self._load_text(file_path)
        else:
            logger.warning(f"Unsupported file format: {file_extension}")
            return {
                'text': '',
                'metadata': {'filename': file_path.name, 'error': 'Unsupported format'},
                'success': False
            }
    
    def _load_pdf(self, file_path: Path) -> Dict[str, any]:
        """Extract text from PDF files"""
        try:
            text_content = []
            metadata = {'filename': file_path.name, 'type': 'pdf'}
            
            # Try pdfplumber first (better for complex PDFs)
            if pdfplumber:
                with pdfplumber.open(file_path) as pdf:
                    metadata['page_count'] = len(pdf.pages)
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_content.append(text)
            # Fallback to pypdf
            elif pypdf:
                reader = PdfReader(file_path)
                metadata['page_count'] = len(reader.pages)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
            else:
                raise ImportError("No PDF library available")
            
            full_text = '\n\n'.join(text_content)
            
            # If no text extracted and OCR is available, try OCR
            if not full_text.strip() and self.use_ocr:
                logger.info(f"No text extracted from {file_path.name}, attempting OCR...")
                # This would require converting PDF to images first
                pass
            
            return {
                'text': full_text,
                'metadata': metadata,
                'success': True
            }
        
        except Exception as e:
            logger.error(f"Error loading PDF {file_path.name}: {str(e)}")
            return {
                'text': '',
                'metadata': {'filename': file_path.name, 'error': str(e)},
                'success': False
            }
    
    def _load_word(self, file_path: Path) -> Dict[str, any]:
        """Extract text from Word documents"""
        try:
            if Document is None:
                raise ImportError("python-docx not installed")
            
            doc = Document(file_path)
            text_content = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_content.append(row_text)
            
            full_text = '\n\n'.join(text_content)
            
            metadata = {
                'filename': file_path.name,
                'type': 'word',
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables)
            }
            
            return {
                'text': full_text,
                'metadata': metadata,
                'success': True
            }
        
        except Exception as e:
            logger.error(f"Error loading Word document {file_path.name}: {str(e)}")
            return {
                'text': '',
                'metadata': {'filename': file_path.name, 'error': str(e)},
                'success': False
            }
    
    def _load_powerpoint(self, file_path: Path) -> Dict[str, any]:
        """Extract text from PowerPoint presentations"""
        try:
            if Presentation is None:
                raise ImportError("python-pptx not installed")
            
            prs = Presentation(file_path)
            text_content = []
            
            # Extract text from each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # Extract text from tables
                    if shape.shape_type == 19:  # Table
                        try:
                            for row in shape.table.rows:
                                row_text = ' | '.join(cell.text for cell in row.cells)
                                if row_text.strip():
                                    slide_text.append(row_text)
                        except:
                            pass
                
                if slide_text:
                    text_content.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_text))
            
            full_text = '\n\n'.join(text_content)
            
            metadata = {
                'filename': file_path.name,
                'type': 'powerpoint',
                'slide_count': len(prs.slides)
            }
            
            return {
                'text': full_text,
                'metadata': metadata,
                'success': True
            }
        
        except Exception as e:
            logger.error(f"Error loading PowerPoint file {file_path.name}: {str(e)}")
            return {
                'text': '',
                'metadata': {'filename': file_path.name, 'error': str(e)},
                'success': False
            }
            return {
                'text': '',
                'metadata': {'filename': file_path.name, 'error': str(e)},
                'success': False
            }
    
    def _load_image(self, file_path: Path) -> Dict[str, any]:
        """Extract text from images using OCR"""
        try:
            if not self.use_ocr or Image is None or pytesseract is None:
                raise ImportError("OCR dependencies not available")
            
            image = Image.open(file_path)
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            
            metadata = {
                'filename': file_path.name,
                'type': 'image',
                'size': image.size,
                'mode': image.mode
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'success': True
            }
        
        except Exception as e:
            logger.error(f"Error loading image {file_path.name}: {str(e)}")
            return {
                'text': '',
                'metadata': {'filename': file_path.name, 'error': str(e)},
                'success': False
            }
    
    def _load_text(self, file_path: Path) -> Dict[str, any]:
        """Load plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            metadata = {
                'filename': file_path.name,
                'type': 'text',
                'size': file_path.stat().st_size
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'success': True
            }
        
        except Exception as e:
            logger.error(f"Error loading text file {file_path.name}: {str(e)}")
            return {
                'text': '',
                'metadata': {'filename': file_path.name, 'error': str(e)},
                'success': False
            }
    
    def load_directory(self, directory_path: str, recursive: bool = True) -> List[Dict[str, any]]:
        """
        Load all supported documents from a directory
        
        Args:
            directory_path: Path to the directory
            recursive: Whether to search subdirectories
            
        Returns:
            List of document dictionaries
        """
        directory = Path(directory_path)
        documents = []
        
        if not directory.exists():
            logger.error(f"Directory not found: {directory}")
            return documents
        
        # Get all supported files
        pattern = '**/*' if recursive else '*'
        
        for file_path in directory.glob(pattern):
            if file_path.is_file():
                ext = file_path.suffix.lower()
                if any(ext in formats for formats in self.SUPPORTED_FORMATS.values()):
                    logger.info(f"Loading: {file_path.name}")
                    doc = self.load_document(str(file_path))
                    if doc['success']:
                        documents.append(doc)
        
        logger.info(f"Loaded {len(documents)} documents from {directory}")
        return documents
